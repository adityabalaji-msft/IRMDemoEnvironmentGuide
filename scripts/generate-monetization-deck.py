"""
Generate IRM & Agent Monetization Strategy Deck (.pptx)
Based on: Monetization strategy - Agent and IRM meeting (July 21, 2026)
Run: pip install python-pptx && python scripts/generate-monetization-deck.py
Output: IRM-Monetization-Strategy.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
LIGHT_BLUE = RGBColor(0x50, 0xE6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xD2, 0xD0, 0xCE)
RED = RGBColor(0xD1, 0x34, 0x38)
GREEN = RGBColor(0x10, 0x7C, 0x10)
AMBER = RGBColor(0xFF, 0xB9, 0x00)
BLACK = RGBColor(0x24, 0x24, 0x24)
ACCENT_TEAL = RGBColor(0x00, 0xB7, 0xC3)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide():
    return prs.slides.add_slide(prs.slide_layouts[5])  # blank


def add_textbox(slide, text, left, top, width, height,
                font_size=Pt(18), color=BLACK, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_multiline(slide, lines, left, top, width, height, font_size=Pt(17), color=BLACK,
                  line_spacing=Pt(8), bold_lines=None):
    """Add multi-line text. bold_lines is a set of 0-based indices to bold."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    bold_lines = bold_lines or set()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.bold = i in bold_lines
        p.space_after = line_spacing
    return tf


def add_note(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def full_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def accent_bar(slide, top, height=Inches(0.06), color=AZURE_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), top, Inches(12.1), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def card_box(slide, left, top, width, height, fill_color=RGBColor(0xF3, 0xF2, 0xF1)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LIGHT_GRAY
    shape.line.width = Pt(1)
    return shape


# ============================================================
# SLIDE 1: Title
# ============================================================
s = add_slide()
full_bg(s, DARK_BLUE)
add_textbox(s, "IRM & Agent", Inches(1), Inches(2.0), Inches(11), Inches(1.2),
            Pt(48), WHITE, True)
add_textbox(s, "Monetization Strategy", Inches(1), Inches(3.0), Inches(11), Inches(1.0),
            Pt(40), LIGHT_BLUE, False)
accent_bar(s, Inches(4.2), color=LIGHT_BLUE)
add_textbox(s, "Proposed Business Model, Principles & Rationale", Inches(1), Inches(4.5),
            Inches(11), Inches(0.6), Pt(20), WHITE)
add_textbox(s, "July 2026  •  Confidential", Inches(1), Inches(6.0),
            Inches(11), Inches(0.5), Pt(14), LIGHT_BLUE)

add_note(s, "Title slide. Set context: This deck captures the monetization strategy aligned "
         "in the July 21 strategy session. Purpose is to document the rationale for LT/Finance.")

# ============================================================
# SLIDE 2: Executive Summary
# ============================================================
s = add_slide()
add_textbox(s, "Executive Summary", Inches(0.6), Inches(0.4), Inches(12), Inches(0.8),
            Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

add_multiline(s, [
    "The team aligned on three key decisions:",
    "",
    "1.  IRM and Agent monetization should be treated separately",
    "     IRM capabilities priced independently from Agent usage, giving flexibility",
    "     to evolve each without blocking the other.",
    "",
    "2.  Proceed with Option B (adoption-first pricing)",
    "     Free/freemium entry point with capability-level paid tiers.",
    "     Avoid reopening the decision — document rationale and move forward.",
    "",
    "3.  Prioritize adoption over margin optimization",
    "     This is a new product category. Revenue and customer value must be",
    "     established before aggressively optimizing profitability.",
], Inches(0.6), Inches(1.6), Inches(12), Inches(5.0), Pt(18),
    bold_lines={0, 2, 7, 11})

add_note(s, "Three headline decisions. Keep this crisp — this is the 'so what' slide.")

# ============================================================
# SLIDE 3: Context — Why This Matters
# ============================================================
s = add_slide()
add_textbox(s, "Context: Why Monetization Strategy Matters Now",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

# Three cards
for i, (title, body) in enumerate([
    ("New Product Category",
     "IRM is establishing a new category\nin Azure — resiliency management\nas a platform service. Customers\nare still learning the value prop."),
    ("New Agent Paradigm",
     "The Resiliency Agent introduces\nan AI-native interaction model.\nAdoption patterns are unproven.\nValue curve is still being defined."),
    ("Finance & LT Alignment",
     "Upcoming discussions with Finance\nrequire a clear, defensible rationale\nfor free-tier justification, pricing,\nand margin expectations."),
]):
    left = Inches(0.6 + i * 4.1)
    card_box(s, left, Inches(1.8), Inches(3.8), Inches(4.5))
    add_textbox(s, title, left + Inches(0.3), Inches(2.0), Inches(3.2), Inches(0.6),
                Pt(20), AZURE_BLUE, True)
    add_textbox(s, body, left + Inches(0.3), Inches(2.7), Inches(3.2), Inches(3.0),
                Pt(16), BLACK)

add_note(s, "Set context: why we need a clear strategy now. Three converging forces.")

# ============================================================
# SLIDE 4: Principle 1 — Adoption Before Margin
# ============================================================
s = add_slide()
add_textbox(s, "Principle 1", Inches(0.6), Inches(0.4), Inches(4), Inches(0.6),
            Pt(16), AZURE_BLUE, True)
add_textbox(s, "Adoption Before Margin", Inches(0.6), Inches(0.8), Inches(12), Inches(0.8),
            Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.5))

add_multiline(s, [
    "Do not optimize for margin first.",
    "",
    "Rationale:",
    "•  This is a new product category — customer awareness is still developing",
    "•  Adoption is more important than maximizing gross margin at this stage",
    "•  Revenue, adoption, and customer value must be established before",
    "    aggressively optimizing profitability",
    "",
    "Messaging to Finance & LT:",
    "•  IRM is a value-added service — not a mature infrastructure commodity",
    "•  Should not be evaluated using the same margin expectations as",
    "    established Azure services (e.g., Compute, Storage)",
    "•  Early-stage adoption is the primary objective",
    "",
    "Reaching 70% margin requires ~30,000 service groups with high conversion",
    "— current service group counts are significantly lower. Margin targets",
    "should be revisited as adoption scales.",
], Inches(0.6), Inches(1.8), Inches(12), Inches(5.0), Pt(17),
    bold_lines={0, 2, 8, 14})

add_note(s, "Strongest message from the meeting. Amit was clear: do not lead with margin. "
         "Lead with adoption. Frame IRM differently from mature infra services.")

# ============================================================
# SLIDE 5: Principle 2 — Stick with Option B
# ============================================================
s = add_slide()
add_textbox(s, "Principle 2", Inches(0.6), Inches(0.4), Inches(4), Inches(0.6),
            Pt(16), AZURE_BLUE, True)
add_textbox(s, "Stick with Option B — Stop Reopening", Inches(0.6), Inches(0.8),
            Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.5))

add_multiline(s, [
    "Several options were evaluated. Option B was selected because:",
    "",
    "•  The product is new — pricing simplicity reduces friction",
    "•  The agent is new — usage patterns are not yet established",
    "•  Customers need exposure to capabilities before paying",
    "•  Adoption is the priority at this stage",
    "",
    "Recommended action:",
    "•  Create a concise artifact documenting:",
    "    – Options considered",
    "    – Pros and cons of each",
    "    – Why Option B was chosen",
    "•  Use this as the reference document for all future discussions",
    "•  Do not reopen the decision — iterate on execution instead",
], Inches(0.6), Inches(1.8), Inches(12), Inches(5.0), Pt(17),
    bold_lines={0, 7})

add_note(s, "Amit emphasized: stop relitigating the decision. Document it, reference it, move on.")

# ============================================================
# SLIDE 6: Principle 3 — Customers Need a Hook
# ============================================================
s = add_slide()
add_textbox(s, "Principle 3", Inches(0.6), Inches(0.4), Inches(4), Inches(0.6),
            Pt(16), AZURE_BLUE, True)
add_textbox(s, 'Customers Need a "Hook"', Inches(0.6), Inches(0.8),
            Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.5))

# Flow diagram via cards
steps = [
    ("1. Show Value\nUpfront", "Free posture assessment,\nat-scale visibility,\nrecommendations", GREEN),
    ("2. Create\nEngagement", "Limited free usage of\npremium features —\ncuriosity & exploration", AMBER),
    ("3. Convert to\nPaid", "Drills, recovery plans,\nadvanced capabilities\nrequire paid tier", AZURE_BLUE),
]
for i, (title, body, color) in enumerate(steps):
    left = Inches(0.6 + i * 4.1)
    shape = card_box(s, left, Inches(2.0), Inches(3.8), Inches(3.0))
    # Color strip at top
    strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, Inches(2.0), Inches(3.8), Inches(0.12))
    strip.fill.solid()
    strip.fill.fore_color.rgb = color
    strip.line.fill.background()
    add_textbox(s, title, left + Inches(0.3), Inches(2.3), Inches(3.2), Inches(1.0),
                Pt(18), DARK_BLUE, True)
    add_textbox(s, body, left + Inches(0.3), Inches(3.5), Inches(3.2), Inches(1.2),
                Pt(15), BLACK)

# Arrow connectors (simple text)
add_textbox(s, "→", Inches(4.2), Inches(3.0), Inches(0.6), Inches(0.6), Pt(36), GRAY, True, PP_ALIGN.CENTER)
add_textbox(s, "→", Inches(8.3), Inches(3.0), Inches(0.6), Inches(0.6), Pt(36), GRAY, True, PP_ALIGN.CENTER)

add_multiline(s, [
    "Analogies discussed: LinkedIn Premium, freemium SaaS, sports channel packages",
    "Key: Show meaningful value before asking customers to pay",
], Inches(0.6), Inches(5.4), Inches(12), Inches(1.5), Pt(16), GRAY)

add_note(s, "The 'hook' concept was central. Give enough value free to prove the product, "
         "then convert. Don't gate everything behind a paywall from day 1.")

# ============================================================
# SLIDE 7: Business Model — IRM vs Agent Separation
# ============================================================
s = add_slide()
add_textbox(s, "Business Model: Separate IRM & Agent Monetization",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

# Two columns
# IRM column
card_box(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.6), Inches(5.8), Inches(0.12))
strip.fill.solid()
strip.fill.fore_color.rgb = AZURE_BLUE
strip.line.fill.background()
add_textbox(s, "IRM (Platform Capabilities)", Inches(0.9), Inches(1.9), Inches(5.2), Inches(0.5),
            Pt(22), AZURE_BLUE, True)
add_multiline(s, [
    "Free tier:",
    "  •  Posture assessment & recommendations",
    "  •  At-scale visibility across service groups",
    "  •  Zone resiliency status dashboard",
    "",
    "Paid capabilities:",
    "  •  Zone-down drills (fault injection)",
    "  •  Recovery orchestration (plans + execution)",
    "",
    "Premium add-ons:",
    "  •  Regional resiliency plans",
    "  •  Dependency discovery",
    "  •  Advanced topology capabilities",
], Inches(0.9), Inches(2.5), Inches(5.2), Inches(4.0), Pt(15))

# Agent column
card_box(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.2))
strip = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.9), Inches(1.6), Inches(5.8), Inches(0.12))
strip.fill.solid()
strip.fill.fore_color.rgb = ACCENT_TEAL
strip.line.fill.background()
add_textbox(s, "Agent (AI-Native Experience)", Inches(7.2), Inches(1.9), Inches(5.2), Inches(0.5),
            Pt(22), ACCENT_TEAL, True)
add_multiline(s, [
    "Pricing model:",
    "  •  Usage-based — tied to agent interactions",
    "  •  Linked to JTBD execution outcomes",
    "",
    "Why separate:",
    "  •  Agent adoption patterns are unproven",
    "  •  Value curve is still being defined",
    "  •  Coupling would block IRM GA pricing",
    "  •  Flexibility to iterate independently",
    "",
    "Open question:",
    "  What capabilities must exist before",
    "  customers prefer the agent over the portal?",
], Inches(7.2), Inches(2.5), Inches(5.2), Inches(4.0), Pt(15))

add_note(s, "Core structural decision: don't couple IRM and Agent pricing. "
         "Each needs room to evolve. Agent value curve is still TBD — that's a product exercise, not pricing.")

# ============================================================
# SLIDE 8: Pricing Philosophy
# ============================================================
s = add_slide()
add_textbox(s, "Pricing Philosophy: Capability-First, Bundles Later",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

# Near-term card
card_box(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(5.2))
add_textbox(s, "Near Term: Capability-Level Pricing", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5),
            Pt(20), AZURE_BLUE, True)
add_multiline(s, [
    "Why individual capabilities first:",
    "",
    "•  Product capabilities are still emerging",
    "•  Customers don't yet understand the full portfolio",
    "•  Individual pricing allows better adoption measurement",
    "•  Enables clearer investment decisions per capability",
    "",
    "Examples of priced capabilities:",
    "  – Zone-down drills",
    "  – Recovery orchestration",
    "  – Regional resiliency planning",
    "  – Dependency discovery",
], Inches(0.9), Inches(2.4), Inches(5.2), Inches(4.0), Pt(16))

# Long-term card
card_box(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(5.2))
add_textbox(s, "Long Term: Evolve Toward Bundles", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5),
            Pt(20), AZURE_BLUE, True)
add_multiline(s, [
    "Once adoption & capabilities mature:",
    "",
    "•  Introduce tiered bundles",
    "•  Create premium packages",
    "•  Offer bundle discounts for commitment",
    "",
    "Analogies:",
    "  – Base package vs. premium package",
    "  – Sports channel bundling models",
    "",
    "Key: Don't bundle before you understand",
    "which capabilities drive the most value.",
    "Let individual pricing data inform bundles.",
], Inches(7.2), Inches(2.4), Inches(5.2), Inches(4.0), Pt(16))

# Arrow between
add_textbox(s, "→", Inches(6.15), Inches(3.6), Inches(0.8), Inches(0.8),
            Pt(40), AZURE_BLUE, True, PP_ALIGN.CENTER)

add_note(s, "Don't jump to bundles. Price capabilities individually to learn what customers value. "
         "Bundle later when you have data.")

# ============================================================
# SLIDE 9: Free / Freemium Tier Design
# ============================================================
s = add_slide()
add_textbox(s, "Free / Freemium Tier: What's Included",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

# Table-like layout
headers = ["Capability", "Free Tier", "Paid Tier"]
rows = [
    ["Posture Assessment", "✅  Included", "—"],
    ["At-Scale Visibility", "✅  Included", "—"],
    ["Recommendations", "✅  Included", "—"],
    ["Zone-Down Drills", "❌  Not included", "✅  Per-drill or subscription"],
    ["Recovery Orchestration", "❌  Not included", "✅  Per-plan or subscription"],
    ["Regional Resiliency", "❌  Not included", "✅  Premium add-on"],
    ["Dependency Discovery", "❌  Not included", "✅  Premium add-on"],
    ["Agent Interactions", "Limited free usage", "✅  Usage-based pricing"],
]

# Header row
for j, h in enumerate(headers):
    left = Inches(0.6 + j * 4.0)
    w = Inches(3.8)
    shape = card_box(s, left, Inches(1.6), w, Inches(0.5), AZURE_BLUE)
    add_textbox(s, h, left + Inches(0.2), Inches(1.62), Inches(3.4), Inches(0.5),
                Pt(15), WHITE, True)

# Data rows
for i, row in enumerate(rows):
    row_top = Inches(2.2 + i * 0.56)
    bg_color = RGBColor(0xF3, 0xF2, 0xF1) if i % 2 == 0 else WHITE
    for j, cell in enumerate(row):
        left = Inches(0.6 + j * 4.0)
        card_box(s, left, row_top, Inches(3.8), Inches(0.5), bg_color)
        c = GREEN if "✅" in cell else (RED if "❌" in cell else BLACK)
        add_textbox(s, cell, left + Inches(0.2), row_top + Inches(0.02), Inches(3.4), Inches(0.5),
                    Pt(14), c)

add_note(s, "Clear delineation: assessment and visibility are free (the hook). "
         "Action-oriented capabilities (drills, recovery, advanced features) are paid.")

# ============================================================
# SLIDE 10: Margin Discussion
# ============================================================
s = add_slide()
add_textbox(s, "Margin Expectations: A Realistic View",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

# Challenge box
card_box(s, Inches(0.6), Inches(1.6), Inches(5.8), Inches(3.0), RGBColor(0xFF, 0xF4, 0xCE))
add_textbox(s, "⚠  The Challenge", Inches(0.9), Inches(1.8), Inches(5.2), Inches(0.5),
            Pt(20), DARK_BLUE, True)
add_multiline(s, [
    "Achieving 70% margin requires aggressive assumptions:",
    "",
    "•  ~30,000 service groups",
    "•  High paid-tier conversion rates",
    "•  Large attach rates for add-ons",
    "",
    "Current service group counts are significantly",
    "lower. Reaching 30K quickly is unrealistic.",
], Inches(0.9), Inches(2.4), Inches(5.2), Inches(2.0), Pt(15))

# Direction box
card_box(s, Inches(6.9), Inches(1.6), Inches(5.8), Inches(3.0), RGBColor(0xDF, 0xF6, 0xDD))
add_textbox(s, "✅  Recommended Direction", Inches(7.2), Inches(1.8), Inches(5.2), Inches(0.5),
            Pt(20), DARK_BLUE, True)
add_multiline(s, [
    "Lead with adoption metrics, not margin targets:",
    "",
    "•  Challenge whether 70-80% is the right target",
    "•  Frame IRM as an early-stage value-added service",
    "•  Margin is a secondary optimization",
    "•  Adoption & customer value come first",
    "",
    "Competitors (e.g., AWS) charge substantially more.",
    "Current ~$5 placeholders may undervalue the offering.",
], Inches(7.2), Inches(2.4), Inches(5.2), Inches(2.0), Pt(15))

add_multiline(s, [
    "Key message for Finance:",
    '"This is a new product in a new category. Evaluate it on adoption trajectory and customer value —',
    ' not on the same margin benchmarks as mature infrastructure services."',
], Inches(0.6), Inches(5.2), Inches(12), Inches(1.8), Pt(18), DARK_BLUE,
    bold_lines={0})

add_note(s, "Be direct with Finance: margin expectations from Compute/Storage don't apply here. "
         "This is early stage. Adoption is the metric that matters.")

# ============================================================
# SLIDE 11: Agent Value Curve
# ============================================================
s = add_slide()
add_textbox(s, "Open Question: Agent Value Curve",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

add_multiline(s, [
    "Before pricing the Agent, we need to answer a fundamental product question:",
], Inches(0.6), Inches(1.6), Inches(12), Inches(0.6), Pt(18))

# Highlighted question
shape = card_box(s, Inches(0.6), Inches(2.3), Inches(12.1), Inches(1.2), RGBColor(0xE8, 0xF0, 0xFE))
add_textbox(s, "What capabilities must exist before customers genuinely prefer\nthe Agent over the portal?",
            Inches(1.0), Inches(2.4), Inches(11.3), Inches(1.0), Pt(24), AZURE_BLUE, True, PP_ALIGN.CENTER)

add_multiline(s, [
    "This is a product strategy exercise, not a monetization exercise.",
    "",
    "Recommended next steps:",
    "",
    '•  Build a customer value curve for the Resiliency Agent',
    '•  Define what "agent delight" looks like — when is the agent genuinely',
    '    better than clicking through the portal?',
    '•  Identify missing capabilities required to achieve meaningful adoption',
    '•  Use these insights to inform agent pricing model (usage-based,',
    '    outcome-based, or hybrid)',
    "",
    "Until this work is done, keep agent pricing separate and flexible.",
], Inches(0.6), Inches(3.8), Inches(12), Inches(3.5), Pt(17),
    bold_lines={0, 2})

add_note(s, "Kirushna raised this: we can't price the agent until we understand its value curve. "
         "This is a product strategy workstream, not a pricing one.")

# ============================================================
# SLIDE 12: Competitive Context
# ============================================================
s = add_slide()
add_textbox(s, "Competitive Context & Pricing Signals",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

add_multiline(s, [
    "Key observations from the discussion:",
    "",
    "•  AWS charges substantially higher for comparable resiliency capabilities",
    "•  Current placeholder prices (~$5 per service group) may significantly",
    "    undervalue the offering",
    "•  Freemium model is common in this space — but premium tiers must",
    "    capture real value",
    "",
    "Pricing should reflect:",
    "",
    "•  The cost of NOT having zone resiliency (outage impact)",
    "•  The value of validated, tested recovery (measured RTO)",
    "•  The operational savings from automated assessment vs. manual review",
    "•  Competitive positioning — don't underprice relative to market",
    "",
    "Action: Review competitive pricing data before finalizing tier pricing.",
], Inches(0.6), Inches(1.6), Inches(12), Inches(5.0), Pt(17),
    bold_lines={0, 8, 15})

add_note(s, "Don't leave money on the table. The value of preventing a zone outage is massive. "
         "Price should reflect that, even with a freemium entry point.")

# ============================================================
# SLIDE 13: Summary — What We Aligned On
# ============================================================
s = add_slide()
add_textbox(s, "Summary: What We Aligned On",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

items = [
    ("1", "Continue with Option B", "Stop reopening. Document rationale and execute."),
    ("2", "Adoption before margin", "New product, new category. Margin is secondary."),
    ("3", "Separate IRM & Agent pricing", "Flexibility to evolve each independently."),
    ("4", "Free/freemium entry point", 'Show value upfront. Convert via the "hook."'),
    ("5", "Capability-level pricing first", "Price individual features. Bundle later with data."),
    ("6", "Build Agent value curve", "Product exercise first, pricing exercise second."),
]

for i, (num, title, desc) in enumerate(items):
    top = Inches(1.5 + i * 0.9)
    # Number circle
    circle = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), top + Inches(0.05), Inches(0.5), Inches(0.5))
    circle.fill.solid()
    circle.fill.fore_color.rgb = AZURE_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    add_textbox(s, title, Inches(1.3), top, Inches(4.5), Inches(0.5), Pt(19), DARK_BLUE, True)
    add_textbox(s, desc, Inches(5.8), top + Inches(0.03), Inches(7), Inches(0.5), Pt(16), GRAY)

add_note(s, "Recap slide. Use this as the anchor when presenting to LT/Finance.")

# ============================================================
# SLIDE 14: Next Steps
# ============================================================
s = add_slide()
add_textbox(s, "Next Steps",
            Inches(0.6), Inches(0.4), Inches(12), Inches(0.8), Pt(32), DARK_BLUE, True)
accent_bar(s, Inches(1.15))

steps = [
    ("Document Option B rationale", "Create concise artifact: options evaluated, pros/cons, why B was chosen", "Aditya"),
    ("Position strategy for LT/Finance", "Frame around: new product, new agent, adoption-first philosophy", "Team"),
    ("Define Agent value curve", "Product exercise: what capabilities = agent delight?", "Kirushna + PM"),
    ("Review competitive pricing data", "Validate tier pricing against AWS and market benchmarks", "Aditya"),
    ("Build free-tier justification", "Prepare materials for Finance on freemium rationale", "Team"),
    ("Finalize capability pricing", "Individual capability prices for drills, recovery, add-ons", "Aditya + Finance"),
]

for i, (action, detail, owner) in enumerate(steps):
    top = Inches(1.5 + i * 0.9)
    card_box(s, Inches(0.6), top, Inches(12.1), Inches(0.75))
    add_textbox(s, action, Inches(0.9), top + Inches(0.05), Inches(4.5), Inches(0.5),
                Pt(16), DARK_BLUE, True)
    add_textbox(s, detail, Inches(5.5), top + Inches(0.05), Inches(5.5), Inches(0.5),
                Pt(14), GRAY)
    add_textbox(s, owner, Inches(11.2), top + Inches(0.05), Inches(1.3), Inches(0.5),
                Pt(14), AZURE_BLUE, True, PP_ALIGN.RIGHT)

add_note(s, "Action items with owners. Review and assign dates in the next sync.")

# ============================================================
# SLIDE 15: One-Line Takeaway
# ============================================================
s = add_slide()
full_bg(s, DARK_BLUE)
add_textbox(s, "One-Line Takeaway", Inches(1), Inches(2.0), Inches(11), Inches(0.8),
            Pt(24), LIGHT_BLUE, False, PP_ALIGN.CENTER)
add_textbox(s, "Adoption-first monetization:", Inches(1), Inches(3.0), Inches(11.3), Inches(0.8),
            Pt(32), WHITE, True, PP_ALIGN.CENTER)
add_multiline(s, [
    "Free/freemium entry  •  IRM priced separately from Agent",
    "Capability-level paid tiers  •  Bundles deferred until maturity",
    "Margin targets revisited after meaningful adoption",
], Inches(1), Inches(4.0), Inches(11.3), Inches(2.5), Pt(22), LIGHT_BLUE)
for p in s.shapes[-1].text_frame.paragraphs:
    p.alignment = PP_ALIGN.CENTER

add_note(s, "Closing slide. Leave this on screen for discussion.")

# ============================================================
# Save
# ============================================================
output_path = "IRM-Monetization-Strategy.pptx"
prs.save(output_path)
print(f"✅ Generated: {output_path}")
print(f"   Slides: {len(prs.slides)}")
