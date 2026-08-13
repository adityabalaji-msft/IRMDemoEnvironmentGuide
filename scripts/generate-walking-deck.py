"""
Generate IRM Field Demo Walking Deck (.pptx)
Run: pip install python-pptx && python scripts/generate-walking-deck.py
Output: IRM-Field-Demo-Walking-Deck.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BLUE = RGBColor(0x00, 0x2B, 0x5C)
LIGHT_BLUE = RGBColor(0x50, 0xE6, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6E, 0x6E, 0x6E)
LIGHT_GRAY = RGBColor(0xF3, 0xF2, 0xF1)
RED = RGBColor(0xD1, 0x34, 0x38)
GREEN = RGBColor(0x10, 0x7C, 0x10)
YELLOW = RGBColor(0xFF, 0xB9, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_slide(title_text, subtitle_text=None, layout_idx=5):
    """Add a blank slide and return it."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    return slide


def add_title_box(slide, text, top=Inches(0.4), left=Inches(0.6), width=Inches(12), height=Inches(1.0),
                  font_size=Pt(32), color=DARK_BLUE, bold=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = color
    p.font.bold = bold
    return txBox


def add_body_text(slide, text, top=Inches(1.6), left=Inches(0.6), width=Inches(12), height=Inches(5.0),
                  font_size=Pt(18), color=BLACK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    lines = text.split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = font_size
        p.font.color.rgb = color
        p.space_after = Pt(6)
        if line.startswith('•'):
            p.level = 1
    return txBox


def add_speaker_note(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_action_banner(slide, text, top=Inches(6.4)):
    """Add a colored banner indicating presenter action (switch to browser, portal, etc.)"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), top, Inches(12.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = AZURE_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = add_slide("", layout_idx=5)
# Background shape
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()

add_title_box(slide, "Infrastructure Resiliency Manager", top=Inches(2.0), left=Inches(1.0),
              font_size=Pt(44), color=WHITE)
add_title_box(slide, "Field Demo Walking Deck", top=Inches(3.0), left=Inches(1.0),
              font_size=Pt(28), color=LIGHT_BLUE, bold=False)
add_title_box(slide, "Contoso Retail — Start Resilient, Get Resilient, Stay Resilient",
              top=Inches(4.2), left=Inches(1.0), font_size=Pt(20), color=WHITE, bold=False)
add_title_box(slide, "Three customer journeys  •  Two live apps  •  Agent-led resiliency",
              top=Inches(5.4), left=Inches(1.0), font_size=Pt(16), color=LIGHT_BLUE, bold=False)

add_speaker_note(slide, """Welcome slide. Introduce yourself and set context:
- "Today I'm going to show you how Azure helps you get resilient across three customer journeys — Start Resilient, Get Resilient, and Stay Resilient."
- "We'll use the Resiliency Agent and Infrastructure Resiliency Manager together, with two live applications deployed in Azure."
""")

# ============================================================
# SLIDE 2: Agenda
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Demo Flow — Three Customer Journeys", font_size=Pt(36))
add_body_text(slide, """Phase 1 — "Start Resilient"
       Resilient by default: Agent generates guidance + IaC templates

Phase 2 — "Get Resilient"
       Find what is critical: Meet the apps → IRM assessment → Copilot remediation

Phase 3 — "Stay Resilient"
       Proactive drift resolution: Config drift detection + compliance drills + recovery orchestration""", font_size=Pt(20))

add_speaker_note(slide, """Agenda overview. Keep this brief — 30 seconds max.
"Three phases, each matching a different customer state: deploying something new, hardening existing apps, or keeping a resilient estate from drifting. The Resiliency Agent and IRM meet you wherever you are."
""")

# ============================================================
# SLIDE 3: Quick Reference
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Quick Reference — Demo Resources")
add_body_text(slide, """AKS App URL:          http://irm-demo-aks.westus2.cloudapp.azure.com
VM App URL:           http://zr-inventory-asr.westus2.cloudapp.azure.com

AKS Service Group:    IRMDemoSG5
VM Service Group:     IRMDemoSG8

AKS Resource Group:   zr-demo-rg-4
VM Resource Group:    zr-demo-vm-asr-rg

Region:               West US 2""", font_size=Pt(18))

add_speaker_note(slide, """Reference slide — keep visible or memorize key URLs.
Pre-open browser tabs: both app URLs + IRM portal + Resiliency Agent.
""")

# ============================================================
# SLIDE 4: Phase 1 Title — Start Resilient
# ============================================================
slide = add_slide("", layout_idx=5)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = AZURE_BLUE
bg.line.fill.background()
add_title_box(slide, "Phase 1", top=Inches(2.5), left=Inches(1.0), font_size=Pt(52), color=WHITE)
add_title_box(slide, '"Start Resilient"', top=Inches(3.5), left=Inches(1.0), font_size=Pt(36), color=WHITE, bold=False)
add_title_box(slide, "Resilient by Default — Agent-Led Template Generation", top=Inches(4.5), left=Inches(1.0),
              font_size=Pt(20), color=LIGHT_BLUE, bold=False)

add_speaker_note(slide, """Transition to Phase 1.
"The first customer journey: Start Resilient. This is the greenfield moment — you're deploying something new, and you want to get the architecture right from the outset."
""")

# ============================================================
# SLIDE 5: Start Resilient — Customer Voice
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, 'Customer Voice — Start Resilient')
add_body_text(slide, """"Help me get started with resiliency by default."

Many organizations deploy from IaC templates that were written
before availability zones even existed.

The result: applications go live without the right resilience
configurations, and costly re-architecture is needed later.

The Resiliency Agent ensures that new deployments start with
the right configuration from day zero — eliminating the need
to retrofit resilience after the fact.
""", font_size=Pt(22), color=DARK_BLUE)

add_speaker_note(slide, """Set the scene: "This is the customer who is about to deploy something new. Maybe they're modernizing a legacy app, maybe it's a net-new workload. They want to get it right from the start."
""")

# ============================================================
# SLIDE 6: Start Resilient — Agent Demo
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Demo: Resiliency Agent — Guidance Report + IaC Generation")
add_body_text(slide, """1. Open the Resiliency Agent (Azure Copilot → Resiliency)

2. Describe your application requirements:
   "I need to deploy an e-commerce app on AKS in West US 2 with
    a SQL database for orders and a storage account for product
    images. Generate zone-resilient Bicep templates."

3. The agent responds with:
   • Guidance report — which services need zone redundancy
   • Modular Bicep templates with zone redundancy baked in:
     - AKS: zone-redundant node pools (zones 1, 2, 3)
     - SQL Database: zone redundancy enabled
     - Storage Account: ZRS (Zone-Redundant Storage)
     - Load Balancer: Standard SKU, zone-redundant frontend
   • Cost implications and trade-offs for each choice

4. Show the generated templates — ready to deploy""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO PORTAL → Azure Copilot → Resiliency Agent")

add_speaker_note(slide, """Open the Resiliency Agent in the Azure portal.
Type the prompt and walk through the agent's response.
Highlight: the customer leaves with deployable IaC, not a list of recommendations.
If time allows, show Terraform generation as well.
"The proof point is that you leave this conversation with deployment-ready code."
""")

# ============================================================
# SLIDE 7: Start Resilient — Key Message
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Key Message — Phase 1")
add_body_text(slide, """"The proof point is that you leave this conversation
with deployable, resilient-by-default infrastructure-
as-code in the tooling you already use — Bicep,
Terraform, or ARM templates.

Not a list of recommendations to figure out later."
""", font_size=Pt(24), color=DARK_BLUE)

add_speaker_note(slide, """Let this land. Then transition:
"But what about the apps you already have running? That's the next journey — Get Resilient."
""")

# ============================================================
# SLIDE 8: Phase 2 Title — Get Resilient
# ============================================================
slide = add_slide("", layout_idx=5)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = AZURE_BLUE
bg.line.fill.background()
add_title_box(slide, "Phase 2", top=Inches(2.5), left=Inches(1.0), font_size=Pt(52), color=WHITE)
add_title_box(slide, '"Get Resilient"', top=Inches(3.5), left=Inches(1.0),
              font_size=Pt(36), color=WHITE, bold=False)
add_title_box(slide, "Find What Is Critical — Assessment, Gaps & Remediation", top=Inches(4.5), left=Inches(1.0),
              font_size=Pt(20), color=LIGHT_BLUE, bold=False)

add_speaker_note(slide, """Transition: "Now let's talk about the biggest installed base — the apps that are already running in Azure. How do you find what's actually critical and close the gaps?"
""")

# ============================================================
# SLIDE 9: App A — AKS Architecture
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Meet the Apps — App A: E-Commerce Platform (AKS)")
add_body_text(slide, """• Frontend pods: product catalog + blob storage for static assets
• Backend pods: order processing via Azure SQL
• Container images pulled from Azure Container Registry
• AKS cluster: 3 nodes across Availability Zones 1, 2, 3
• Standard Load Balancer routes traffic across zones

Zone Resiliency Status:
  ✅  AKS Cluster (zones 1/2/3)         ✅  Load Balancer (Standard)
  ✅  Container Registry (zone-redundant)
  ❌  Azure SQL Database (GP_Gen5_2 — no ZR)
  ❌  Storage Account (Standard_LRS — no ZR)""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO BROWSER → Open http://irm-demo-aks.westus2.cloudapp.azure.com")

add_speaker_note(slide, """Open the AKS app in browser. Show it's live.
Key point: "The compute layer looks resilient — nodes across 3 zones. But look at the dependencies:
SQL and Storage are NOT zone-redundant. If Zone 1 goes down, compute survives but data might not."
This is a typical brownfield reality.
""")

# ============================================================
# SLIDE 10: App B — VM Architecture
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Meet the Apps — App B: Inventory Management (VM + ASR)")
add_body_text(slide, """• Monolithic Node.js app on a zone-pinned VM (Zone 1)
• Companion worker VM: data sync agent (also Zone 1)
• Azure Site Recovery: zonal DR (Zone 1 → Zone 2)
• Standard Public IP: zone-redundant by default

Zone Resiliency Status:
  ✅  Main VM (Zone 1, ASR replicates to Zone 2)
  ✅  Worker VM (Zone 1, ASR replicates to Zone 2)
  ✅  Public IP (Standard SKU — zone-redundant)
  ✅  ASR Vault (orchestrates zonal failover)
  ❌  Azure SQL Database (GP_Gen5_2 — no ZR)
  ❌  Storage Account (Standard_LRS — no ZR)""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO BROWSER → Open http://zr-inventory-asr.westus2.cloudapp.azure.com")

add_speaker_note(slide, """Open the VM app in browser. Show it's live.
Key point: "ASR is configured — the VMs CAN recover to Zone 2. But has this ever been tested?
Can Contoso orchestrate recovery in the right order under real pressure?
These are brownfield realities — and this is where IRM helps you get resilient."
""")

# ============================================================
# SLIDE 11: Get Resilient — At-Scale Assessment
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Assess Posture at Scale — IRM Overview")
add_body_text(slide, """Navigate to:  Resiliency → Resiliency Overview

What to show:
• Zone-resilient vs. non-resilient service groups (summary tiles)
• Total resource count broken down by posture
• Color-coded health across all service groups

Talking point:
"This is what a platform team sees when managing dozens of
applications — which apps meet zone resilience goals and
which don't. No manual inspection required."

Pre-created service groups:
  • IRMDemoSG5 → AKS app (Goals + Drill)
  • IRMDemoSG8 → VM app (Goals + Recovery Plan + Drill)""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO PORTAL → IRM: Resiliency Overview (at-scale view)")

add_speaker_note(slide, """Switch to portal now. Navigate to Resiliency Overview.
Point out the non-resilient service groups tile.
"Imagine you manage 50 applications — you need this view to know where to act."
""")

# ============================================================
# SLIDE 12: Get Resilient — Drill into Service Groups
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Drill into Service Groups — Per-Resource Breakdown")
add_body_text(slide, """IRMDemoSG5 (AKS App):
  ✅  AKS Cluster, Load Balancer → Zone-resilient
  ❌  SQL Database, Storage Account → Non zone-resilient
  Recommendations auto-generated with cost indicators

IRMDemoSG8 (VM App):
  ✅  VMs with ASR configured → Zone-resilient
  ❌  SQL Database, Storage → Non zone-resilient
  Recovery Plan already associated for orchestrated failover

For each non-resilient resource, show:
  • Step-by-step remediation guidance
  • Qualitative cost indicator (Low / Medium / High)""", font_size=Pt(17))
add_action_banner(slide, "▶  STAY IN PORTAL → Click into IRMDemoSG5, then IRMDemoSG8")

add_speaker_note(slide, """Walk through both service groups.
Point to green checkmarks and red X's for each resource.
Highlight that VM app has a recovery plan — "IRM understands orchestration needs."
""")

# ============================================================
# SLIDE 13: Get Resilient — Copilot Remediation
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, 'Close the Gaps — Copilot "Resolve" + IaC Generation')
add_body_text(slide, """1. Select a SQL Database recommendation → click "Resolve"
2. Copilot guides the user step by step:

  • What can be fixed in place
    (e.g., portal toggle to enable zone redundancy)

  • What needs to be redeployed via script or automation
    (e.g., storage account LRS → ZRS conversion)

  • What requires manual effort
    (e.g., architecture changes, support requests)

3. Prompt the agent to generate an IaC template (Bicep) with
   zone-redundancy enabled — ready to deploy or integrate
   into existing CI/CD pipelines.""", font_size=Pt(17))
add_action_banner(slide, "▶  STAY IN PORTAL → Click 'Resolve' on a recommendation")

add_speaker_note(slide, """Click Resolve on the SQL Database recommendation.
Walk through the Copilot response — show the categorization of fixes.
If time allows, prompt: "Generate a Bicep template with zone redundancy enabled."
"The agent categorizes each fix by effort and generates deployment-ready code."
""")

# ============================================================
# SLIDE 14: Get Resilient — Key Message
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Key Message — Phase 2")
add_body_text(slide, """"For brownfield estates, this is how you close the gap
between your current posture and your resiliency goal.

IRM doesn't just tell you what's wrong — the agent
categorizes each fix by effort, generates deployment-
ready IaC templates, and ranks the top actions worth
acting on."
""", font_size=Pt(24), color=DARK_BLUE)

add_speaker_note(slide, """Pause — let value land.
Transition: "Assessment done. Gaps closed. But how do we make sure it stays that way?"
""")

# ============================================================
# SLIDE 15: Phase 3 Title — Stay Resilient
# ============================================================
slide = add_slide("", layout_idx=5)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = AZURE_BLUE
bg.line.fill.background()
add_title_box(slide, "Phase 3", top=Inches(2.5), left=Inches(1.0), font_size=Pt(52), color=WHITE)
add_title_box(slide, '"Stay Resilient"', top=Inches(3.5), left=Inches(1.0),
              font_size=Pt(36), color=WHITE, bold=False)
add_title_box(slide, "Proactive Drift Resolution — Drift Detection, Drills & Recovery", top=Inches(4.5), left=Inches(1.0),
              font_size=Pt(20), color=LIGHT_BLUE, bold=False)

add_speaker_note(slide, """Transition: "You've invested in making your apps zone-resilient. But environments drift — new deployments from old templates, configuration changes. How do you keep your estate from silently regressing?"
""")

# ============================================================
# SLIDE 16: Stay Resilient — Config Drift Detection
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Detect Configuration Drift")
add_body_text(slide, """The problem:
  "New deployment, old template — am I still resilient?"

A new deployment from an outdated template can silently regress
your posture. IRM catches this automatically.

What to show:
  1. Navigate to a previously assessed service group
  2. New deployments/changes are automatically reassessed
  3. A resource that was zone-resilient now shows non-compliant
  4. Recommendations guide approved corrections
  5. The agent can apply course correction immediately

Two hooks for the Field:
  • Drift the customer did not know about
  • The regulatory calendar — audits require proof""", font_size=Pt(17))
add_action_banner(slide, "▶  STAY IN PORTAL → Show drift detection in service group posture view")

add_speaker_note(slide, """Show a service group where posture has changed.
Point out the non-compliant resource and the recommendation.
"This is the drift you didn't know about. IRM caught it before a real outage did."
""")

# ============================================================
# SLIDE 17: Stay Resilient — CSA Pre-Demo Drill Option
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "CSA Option: Execute AKS Drill Before the Demo")
add_body_text(slide, """Recommended: Run the AKS drill BEFORE the customer meeting

Why?
  • Avoids live wait times during the demo (drills take minutes)
  • The web app's Activity Log captures the full drill timeline
  • You walk into the meeting with compliance evidence ready

Steps:
  1. Navigate to IRMDemoSG5 → Resiliency → Drills
  2. Execute the drill targeting Zone 1
  3. Wait for completion (~5-10 min)
  4. Open the web app → click "Infrastructure" button
  5. The Activity Log shows: zone down → failover → recovery

During the demo:
  • Open the app and switch to Infrastructure panel
  • Scroll through the Activity Log as audit evidence

Note: The Activity Log uses browser localStorage — use the same
browser session you used to observe the drill.""", font_size=Pt(16))
add_action_banner(slide, "▶  PRE-DEMO: Execute drill, then show Activity Log during customer meeting")

add_speaker_note(slide, """This is the recommended approach for time-constrained demos.
Run the drill 30 minutes before the meeting. The Activity Log persists in the browser.
""")

# ============================================================
# SLIDE 18: Stay Resilient — AKS Drill
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Compliance Drill: AKS App (IRMDemoSG5)")
add_body_text(slide, """Goal: Prove AKS compute survives a zone failure — generate
compliance evidence

Navigate to: IRMDemoSG5 → Resiliency → Drills (pre-created)

Fault Designer shows:
  ✅  AKS Cluster — included for fault injection
       (node shutdown in target zone via Chaos Studio)
  ⛔  SQL Database — excluded from drill
       (non-ZR, would cause expected failures)

Execute drill targeting Zone 1:
  1. AKS node pool VMs in Zone 1 shut down
  2. Load Balancer routes traffic to zones 2 and 3
  3. Open app URL — it continues serving ✅
  4. Show Activity Log — timestamped audit trail
  5. End drill — nodes return, pods rebalance""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO PORTAL → IRMDemoSG5 → Drills → Execute")

add_speaker_note(slide, """Execute the drill or show pre-recorded Activity Log.
"We validated what should survive. The Activity Log provides compliance evidence for audit.
Next step: make SQL zone-redundant, then run the full drill with all resources in scope."
""")

# ============================================================
# SLIDE 19: Stay Resilient — VM Drill
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Recovery Orchestration: VM App (IRMDemoSG8)")
add_body_text(slide, """Goal: Prove the orchestrated recovery plan brings the app
back in the correct sequence

Execute drill targeting Zone 1:
  → Both VMs shut down → App goes DARK ❌

Execute the Recovery Plan (orchestrated sequence):
  1. Worker VM fails over to Zone 2 (data sync agent ready)
  2. Main App VM fails over to Zone 2 (app depends on worker)
  3. ASR handles disk replication, IP reassignment, boot

Validate recovery:
  • App comes back on Zone 2 ✅
  • Check health endpoint for SQL + Storage connectivity
  • Measure RTO from drill metrics

After validation:
  • Reprotect (Zone 2 → Zone 1) for future drills""", font_size=Pt(17))
add_action_banner(slide, "▶  SWITCH TO PORTAL → IRMDemoSG8 → Drills → Recovery Plan")

add_speaker_note(slide, """This is the dramatic moment — the app goes down and comes back.
Show the recovery plan executing in portal.
Switch to browser once recovery completes to show app is live again.
"ASR gives you the capability to recover. The recovery plan proves you can do it in the right order.
Drill metrics give you your measured RTO — compliance evidence for audit."
""")

# ============================================================
# SLIDE 20: Stay Resilient — Key Message
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Key Message — Phase 3")
add_body_text(slide, """"Environments drift. New deployments from old templates
silently regress your posture. IRM catches this
automatically.

Compliance drills generate the evidence auditors need.
Recovery orchestration proves your team can recover
in the right order, under pressure.

That's how you stay resilient."
""", font_size=Pt(24), color=DARK_BLUE)

add_speaker_note(slide, """Let this land. Then move to summary.
""")

# ============================================================
# SLIDE 21: Summary — Three Journeys
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Three Customer Journeys — Summary")
add_body_text(slide, """Phase              Customer Moment            What We Showed                  Outcome
───────────────────────────────────────────────────────────────────────────────────────────
Start Resilient    Greenfield — deploying      Resiliency Agent generates       Deploys resilient
                   something new               guidance + IaC templates         from day zero

Get Resilient      Brownfield — hardening      IRM at-scale assessment →        Knows posture,
                   existing estate             drill-down → Copilot remediation closes gaps

Stay Resilient     Steady-state — keeping      Config drift detection +         Catches regression,
                   what's resilient            compliance drills + recovery     proves readiness
                   from drifting               orchestration                    for audits""", font_size=Pt(16))

add_speaker_note(slide, """Summary slide — quick recap.
"Three journeys. Start right, find what's critical, keep it from drifting.
The Resiliency Agent and IRM meet the customer wherever they are."
""")

# ============================================================
# SLIDE 22: Call to Action
# ============================================================
slide = add_slide("", layout_idx=5)
bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BLUE
bg.line.fill.background()
add_title_box(slide, "Next Steps", top=Inches(2.0), left=Inches(1.0), font_size=Pt(40), color=WHITE)
add_body_text(slide, """• Use the Resiliency Agent to generate resilient templates for new deployments

• Identify 2-3 critical applications to onboard as Service Groups

• Run an assessment — see your zone resiliency posture today

• Create your first drill — prove what works, find what doesn't""",
              top=Inches(3.2), left=Inches(1.0), font_size=Pt(22), color=WHITE)
add_title_box(slide, "Contact: azureresiliency@microsoft.com", top=Inches(6.2), left=Inches(1.0),
              font_size=Pt(18), color=LIGHT_BLUE, bold=False)

add_speaker_note(slide, """Close with actionable next steps.
Start with the agent (Start Resilient) — lowest friction entry point.
Then offer to help set up service groups and run their first drill.
""")

# ============================================================
# SLIDE 23: Appendix — Pre-Demo Checklist
# ============================================================
slide = add_slide("", layout_idx=5)
add_title_box(slide, "Appendix: Pre-Demo Checklist")
add_body_text(slide, """Before presenting, verify:

  □  AKS app is live: http://irm-demo-aks.westus2.cloudapp.azure.com
  □  VM app is live: http://zr-inventory-asr.westus2.cloudapp.azure.com
  □  IRM portal loads and shows IRMDemoSG5 + IRMDemoSG8
  □  Resiliency Agent is accessible (Azure Copilot → Resiliency)
  □  Pre-open browser tabs: both app URLs + IRM portal + Copilot
  □  ASR replication is healthy (portal → Recovery Services vault)
  □  Drill definitions exist in both service groups

Optional — Execute AKS drill ahead of time:
  □  Run drill on IRMDemoSG5 (Zone 1 fault) ~30 min before demo
  □  Keep the AKS app browser tab open during drill execution
  □  Verify Activity Log captured the drill events (Infrastructure panel)
  □  During demo: open Infrastructure → show Activity Log as evidence

If AKS app is down:
  • Check: az aks get-credentials + kubectl get pods
  • Pods may need restart after a previous drill

If VM app is down:
  • Check: az vm run-command (see setup-readme.md)
  • May need to restart the scenario6 systemd service""", font_size=Pt(16))

add_speaker_note(slide, """Hidden slide for presenter prep only. Don't show this during the demo.
Run through this checklist 30 minutes before your session.
The AKS drill pre-execution is RECOMMENDED — the Activity Log in the web app
provides compelling proof of zone resilience without waiting during the live demo.""")

# Save
output_path = "IRM-Field-Demo-Walking-Deck.pptx"
prs.save(output_path)
print(f"✅ Generated: {output_path}")
print(f"   Slides: {len(prs.slides)}")
